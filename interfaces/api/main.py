import json
import asyncio
import subprocess
from pathlib import Path
from fastapi import FastAPI, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import StreamingResponse, FileResponse

from interfaces.api.routes import projects, tasks, articles, analytics, kb_routes, settings_routes, skills_routes

app = FastAPI(title="SEO AI Agent", version="0.1.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Register route modules
app.include_router(projects.router)
app.include_router(tasks.router)
app.include_router(articles.router)
app.include_router(analytics.router)
app.include_router(kb_routes.router)
app.include_router(settings_routes.router)
app.include_router(skills_routes.router)

# Mount web UI static files
web_dir = Path(__file__).parent.parent / "web"
web_dir.mkdir(parents=True, exist_ok=True)
app.mount("/static", StaticFiles(directory=str(web_dir)), name="static")


@app.get("/")
async def index():
    """Serve the web UI"""
    index_path = web_dir / "index.html"
    if index_path.exists():
        return FileResponse(str(index_path))
    return {"message": "SEO AI Agent API", "docs": "/docs", "web_ui": "coming soon"}


@app.get("/health")
async def health():
    from agent.tool_registry import registry
    return {
        "status": "ok",
        "version": "0.1.0",
        "tools_count": len(registry.list_tools()),
    }


@app.post("/api/agent/run")
async def agent_run(request: Request):
    """Run an agent task with SSE streaming of progress."""
    body = await request.json()
    task = body.get("task", "")
    project_id = body.get("project_id", "default")

    if not task:
        return {"error": "task is required"}

    async def event_stream():
        # ── Skill command routing ──────────────────────
        if task.startswith("/pptx"):
            yield f"data: {json.dumps({'type': 'start', 'task': task, 'skill': 'pptx'})}\n\n"
            result = await _handle_pptx_command(task.replace("/pptx", "", 1).strip())
            yield f"data: {json.dumps({'type': 'done', 'content': result})}\n\n"
            yield "data: [DONE]\n\n"
            return
        elif task.startswith("/excel"):
            yield f"data: {json.dumps({'type': 'start', 'task': task, 'skill': 'excel'})}\n\n"
            result = await _handle_excel_command(task.replace("/excel", "", 1).strip())
            yield f"data: {json.dumps({'type': 'done', 'content': result})}\n\n"
            yield "data: [DONE]\n\n"
            return
        elif task.startswith("/browser"):
            yield f"data: {json.dumps({'type': 'start', 'task': task, 'skill': 'browser'})}\n\n"
            result = await _handle_browser_command(task.replace("/browser", "", 1).strip())
            yield f"data: {json.dumps({'type': 'done', 'content': result})}\n\n"
            yield "data: [DONE]\n\n"
            return

        from agent.orchestrator import SEOAgent, AgentContext
        from agent.system_prompt import build_system_prompt
        from llm.types import Message
        from config.settings import settings

        agent = SEOAgent()
        await agent._init()
        ctx = AgentContext(project_id=project_id)

        # Build system prompt
        from memory.user_profile import get_profile
        from memory.structured.keyword_memory import get_all_keywords
        from memory.structured.content_memory import get_recent_articles
        from memory.structured.step_memory import get_token_usage

        persona = await get_profile(ctx.user_id)
        articles = await get_recent_articles(ctx.project_id)
        keywords = await get_all_keywords(ctx.project_id)
        usage = await get_token_usage(ctx.project_id)

        memory_data = {
            "recent_articles": ", ".join(a.get("title", "") for a in articles[:5]) if articles else "none",
            "tracked_keywords": ", ".join(k.get("keyword", "") for k in keywords[:10]) if keywords else "none",
            "total_steps": usage.get("total_steps", 0),
            "total_tokens": usage.get("total_tokens", 0),
        }

        # Auto-RAG
        kb_context = ""
        kb_results = await agent.kb.search(task, project_id=ctx.project_id, top_k=settings.kb_default_top_k)
        if kb_results:
            lines = []
            for r in kb_results:
                src = r.get("metadata", {}).get("filename", "unknown")
                content = r.get("content", "")[:400]
                lines.append(f"  [{src}] {content}")
            kb_context = "\n".join(lines)

        system = build_system_prompt(persona=persona, memory=memory_data, kb_context=kb_context)

        messages: list[Message] = [Message(role="user", content=task)]
        max_iterations = 25
        iteration = 0

        yield f"data: {json.dumps({'type': 'start', 'task': task[:100], 'kb_context': kb_context[:200]})}\n\n"

        while iteration < max_iterations:
            iteration += 1
            tools = agent.registry.list_tools()

            response = await agent.llm.chat(messages=messages, tools=tools, system=system)

            if response.tool_calls:
                yield f"data: {json.dumps({'type': 'thinking', 'content': response.content[:300] if response.content else '', 'tools': [tc.name for tc in response.tool_calls]})}\n\n"

            if not response.tool_calls:
                messages.append(Message(role="assistant", content=response.content))
                if len(response.content) > 100:
                    from tools.skills.export_utils import save_all_formats
                    save_all_formats(response.content, prefix="final-report")
                yield f"data: {json.dumps({'type': 'done', 'content': response.content})}\n\n"
                yield "data: [DONE]\n\n"
                return

            messages.append(Message(role="assistant", content=response.content, tool_calls=response.tool_calls))

            for tc in response.tool_calls:
                result_text = await agent.registry.execute(tc.name, tc.args)
                messages.append(Message(role="tool", content=result_text[:4000], tool_call_id=tc.id))
                yield f"data: {json.dumps({'type': 'tool_result', 'tool': tc.name, 'result': result_text[:300]})}\n\n"

                # Auto-ingest + save as downloadable file
                if tc.name in ("competitor_audit", "serp_analyzer", "keyword_research",
                               "rank_tracker", "report_generator", "copywriter",
                               "outline_generator", "seo_scorer", "readability",
                               "fact_checker", "internal_linker", "schema_markup"):
                    if len(result_text) > 200 and "Error" not in result_text:
                        await agent.kb.ingest_text(result_text, source="tool_output", filename=f"{tc.name}_{ctx.task_id}", project_id=project_id)
                        from tools.skills.export_utils import save_all_formats
                        save_all_formats(result_text, prefix=tc.name)

                from memory.structured.step_memory import log_step
                await log_step(
                    project_id=ctx.project_id, task_id=ctx.task_id,
                    step_type="tool_call", tool_name=tc.name,
                    input_summary=json.dumps(tc.args, ensure_ascii=False),
                    output_summary=result_text[:500],
                    success="Error" not in result_text,
                )

        # Save final output
        final_text = messages[-1].content if messages else ""
        if len(final_text) > 100:
            from tools.skills.export_utils import save_all_formats
            save_all_formats(final_text, prefix="final-report")

        yield f"data: {json.dumps({'type': 'error', 'content': 'Max iterations reached'})}\n\n"
        yield "data: [DONE]\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")


# ── Skill command handlers ──────────────────────────────

async def _handle_pptx_command(content: str) -> str:
    """Handle /pptx command — generate a PowerPoint from Markdown content."""
    if not content:
        return "请提供 PPT 内容。用法: /pptx <Markdown 内容>，用 --- 分隔幻灯片。"

    title = content.split("\n")[0].replace("# ", "").strip()[:80]
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slides_text = content.split("\n---\n")
        slide_count = 0

        for slide_md in slides_text:
            lines = slide_md.strip().split("\n")
            title_text = ""
            body_lines = []

            for line in lines:
                if line.startswith("# ") and not title_text:
                    title_text = line[2:]
                elif line.startswith("## ") and not title_text:
                    title_text = line[3:]
                else:
                    body_lines.append(line)

            if not title_text:
                body_lines = [l for l in lines if l.strip()]
                if body_lines:
                    title_text = body_lines[0][:100]
                    body_lines = body_lines[1:]

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_text[:100]

            if body_lines:
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                for bl in body_lines[:20]:
                    p = tf.add_paragraph()
                    p.text = bl[:200]
                    p.level = 0

            slide_count += 1

        out_dir = Path("data/exports")
        out_dir.mkdir(parents=True, exist_ok=True)
        filename = title.replace(" ", "-").replace("/", "-")[:40]
        out_path = out_dir / f"{filename}.pptx"
        prs.save(str(out_path))

        return f"<span class='tool-badge'>pptx ✓</span> PPT 已生成: **{out_path.name}**（{slide_count} 页）\n\n[点击下载](/api/skills/files/{out_path.name})"
    except ImportError:
        return "<span style='color:#cc3333'>python-pptx 未安装。pip install python-pptx</span>"
    except Exception as e:
        return f"<span style='color:#cc3333'>PPT 生成失败: {e}</span>"


async def _handle_excel_command(content: str) -> str:
    """Handle /excel command — create Excel from JSON or natural language."""
    import json as _json

    try:
        import openpyxl
    except ImportError:
        return "<span style='color:#cc3333'>openpyxl 未安装。pip install openpyxl</span>"

    data = None

    # Try JSON first
    if content.strip():
        try:
            data = _json.loads(content)
            if isinstance(data, dict):
                data = [data]  # single row → wrap in list
        except _json.JSONDecodeError:
            pass

    # If not JSON, try extracting a table from natural language via LLM
    if data is None and content.strip():
        try:
            from llm.deepseek_client import DeepSeekClient
            from llm.types import Message
            client = DeepSeekClient()
            prompt = (
                "Convert the following description into a JSON array of objects for an Excel spreadsheet. "
                "Each object represents one row with column names as keys. Return ONLY valid JSON, no explanation.\n\n"
                f"Description: {content[:2000]}\n\n"
                "Example output: [{\"Column1\":\"value1\",\"Column2\":\"value2\"}]"
            )
            response = await client.chat(messages=[Message(role="user", content=prompt)])
            # Extract JSON from response
            text = response.content.strip()
            if "```" in text:
                text = text.split("```")[1]
                if text.startswith("json"):
                    text = text[4:]
            data = _json.loads(text)
        except Exception:
            pass

    if not data:
        return (
            "<span style='color:#cc3333'>无法解析数据。</span>\n\n"
            "支持两种格式：\n"
            "1. **JSON 数组**: `/excel [{\"关键词\":\"test\",\"搜索量\":1000}]`\n"
            "2. **自然语言**: `/excel 创建一个关键词表格，包含关键词、搜索量、竞争度三列，数据是...`"
        )

    if not isinstance(data, list):
        return "数据必须是数组格式。"

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"

    if isinstance(data[0], dict):
        headers = list(data[0].keys())
        for col, h in enumerate(headers, 1):
            ws.cell(row=1, column=col, value=h)
        for row_idx, row_data in enumerate(data, 2):
            for col_idx, key in enumerate(headers, 1):
                ws.cell(row=row_idx, column=col_idx, value=row_data.get(key, ""))
    elif isinstance(data[0], list):
        for row_idx, row_data in enumerate(data, 1):
            for col_idx, val in enumerate(row_data, 1):
                ws.cell(row=row_idx, column=col_idx, value=val)
    else:
        return "数据格式不支持。每行必须是 dict 或 list。"

    out_dir = Path("data/exports")
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "export.xlsx"
    wb.save(str(out_path))

    return f"<span class='tool-badge'>excel ✓</span> 文件已生成: **{out_path.name}**（{ws.max_row} 行 × {ws.max_column} 列）\n\n[点击下载](/api/skills/files/{out_path.name})"


async def _handle_browser_command(content: str) -> str:
    """Handle /browser command — screenshot a URL."""
    if not content.strip():
        return "请提供 URL。用法: /browser https://example.com"

    url = content.strip().split()[0]
    if not url.startswith("http"):
        url = "https://" + url

    try:
        result = subprocess.run(
            ["agent-browser", "--version"],
            capture_output=True, text=True, timeout=5,
        )
        has_browser = result.returncode == 0
    except Exception:
        has_browser = False

    if not has_browser:
        return "<span class='tool-badge'>browser</span> agent-browser 未安装。安装: `npm i -g agent-browser && agent-browser install`"

    out_dir = Path("data/exports")
    out_dir.mkdir(parents=True, exist_ok=True)
    out_file = out_dir / "screenshot.png"

    try:
        result = subprocess.run(
            ["agent-browser", "screenshot", url, "--output", str(out_file)],
            capture_output=True, text=True, timeout=30,
        )
        if out_file.exists():
            return f"<span class='tool-badge'>browser ✓</span> 截图已保存: **{out_file.name}**\n\n[点击下载](/api/skills/files/{out_file.name})"
        return f"<span style='color:#cc3333'>截图失败: {result.stderr[:200]}</span>"
    except Exception as e:
        return f"<span style='color:#cc3333'>浏览器操作失败: {e}</span>"
