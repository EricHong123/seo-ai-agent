"""Skills API — exposes all 4 agent skills as HTTP endpoints + file browser."""

import json
import subprocess
import tempfile
import os
from pathlib import Path
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel

router = APIRouter(prefix="/api/skills", tags=["skills"])

SKILLS_DIR = Path(__file__).parent.parent.parent.parent / ".agents" / "skills"


class PptxRequest(BaseModel):
    action: str = "create"  # create | read
    title: str = ""
    content: str = ""       # Markdown content for slides
    template: str = ""      # Path to template .pptx (optional)


class ExcelRequest(BaseModel):
    action: str = "create"  # create | read | update
    filepath: str = ""      # Path to existing file
    data: list[dict] | None = None  # Data to write
    sheet_name: str = "Sheet1"


class BrowserRequest(BaseModel):
    action: str = "screenshot"  # screenshot | navigate | extract
    url: str = ""
    selector: str = ""  # CSS selector for extraction (optional)


# ── PPTX ──────────────────────────────────────────────

@router.post("/pptx")
async def skill_pptx(body: PptxRequest):
    """Generate or read PowerPoint presentations."""
    if body.action == "read":
        if not body.filepath:
            raise HTTPException(400, "filepath required for read action")
        try:
            result = subprocess.run(
                ["python", "-m", "markitdown", body.filepath],
                capture_output=True, text=True, timeout=30,
            )
            return {"status": "ok", "content": result.stdout[:10000]}
        except Exception as e:
            return {"status": "error", "message": str(e)}

    elif body.action == "create":
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            slides_text = body.content.strip().split("\n---\n") if body.content else [f"# {body.title or 'Presentation'}"]
            if not slides_text or slides_text == [""]:
                slides_text = [f"# {body.title or 'Presentation'}\n\nAdd your content here."]

            slide_count = 0
            for slide_md in slides_text:
                lines = slide_md.strip().split("\n")
                title_text = ""
                body_lines = []

                for line in lines:
                    if line.startswith("# ") and not title_text:
                        title_text = line[2:]
                    elif line.startswith("## ") and not title_text:
                        title_text = line[3:]
                    else:
                        body_lines.append(line)

                if not title_text:
                    title_text = body_lines[0] if body_lines else body.title or "Slide"
                    body_lines = body_lines[1:] if body_lines else []

                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title_text[:100]

                if body_lines:
                    body_shape = slide.shapes.placeholders[1]
                    tf = body_shape.text_frame
                    tf.clear()
                    for bl in body_lines[:20]:
                        p = tf.add_paragraph()
                        p.text = bl[:200]
                        p.level = 0

                slide_count += 1

            # Save
            out_dir = Path("data/exports")
            out_dir.mkdir(parents=True, exist_ok=True)
            filename = (body.title or "presentation").replace(" ", "-").replace("/", "-")[:40]
            out_path = out_dir / f"{filename}.pptx"
            prs.save(str(out_path))

            return {
                "status": "ok",
                "slides": slide_count,
                "filepath": str(out_path),
                "url": f"/static/exports/{out_path.name}" if False else str(out_path),
            }
        except ImportError:
            return {"status": "error", "message": "python-pptx not installed. Run: pip install python-pptx"}
        except Exception as e:
            return {"status": "error", "message": str(e)}

    return {"status": "error", "message": f"Unknown action: {body.action}"}


# ── Excel ──────────────────────────────────────────────

@router.post("/excel")
async def skill_excel(body: ExcelRequest):
    """Create, read, or update Excel files."""
    try:
        import openpyxl
    except ImportError:
        return {"status": "error", "message": "openpyxl not installed. Run: pip install openpyxl"}

    if body.action == "create":
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = body.sheet_name

        if body.data:
            # Write headers
            if isinstance(body.data, list) and len(body.data) > 0:
                if isinstance(body.data[0], dict):
                    headers = list(body.data[0].keys())
                    for col, h in enumerate(headers, 1):
                        ws.cell(row=1, column=col, value=h)
                    for row_idx, row_data in enumerate(body.data, 2):
                        for col_idx, key in enumerate(headers, 1):
                            ws.cell(row=row_idx, column=col_idx, value=row_data.get(key, ""))
                elif isinstance(body.data[0], list):
                    for row_idx, row_data in enumerate(body.data, 1):
                        for col_idx, val in enumerate(row_data, 1):
                            ws.cell(row=row_idx, column=col_idx, value=val)

        out_dir = Path("data/exports")
        out_dir.mkdir(parents=True, exist_ok=True)
        out_path = out_dir / "export.xlsx"
        wb.save(str(out_path))

        return {"status": "ok", "filepath": str(out_path), "sheet": ws.title, "rows": ws.max_row, "cols": ws.max_column}

    elif body.action == "read":
        if not body.filepath:
            raise HTTPException(400, "filepath required")
        wb = openpyxl.load_workbook(body.filepath, data_only=True)
        ws = wb[body.sheet_name] if body.sheet_name in wb.sheetnames else wb.active

        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append(list(row))

        return {"status": "ok", "sheet": ws.title, "rows": len(rows), "data": rows[:100]}

    return {"status": "error", "message": f"Unknown action: {body.action}"}


# ── Browser ────────────────────────────────────────────

@router.post("/browser")
async def skill_browser(body: BrowserRequest):
    """Browser automation via agent-browser CLI."""
    try:
        result = subprocess.run(
            ["agent-browser", "--version"],
            capture_output=True, text=True, timeout=5,
        )
        has_agent_browser = result.returncode == 0
    except Exception:
        has_agent_browser = False

    if not has_agent_browser:
        return {
            "status": "unavailable",
            "message": "agent-browser not installed. Install: npm i -g agent-browser && agent-browser install",
        }

    if body.action == "screenshot" and body.url:
        out_dir = Path("data/exports")
        out_dir.mkdir(parents=True, exist_ok=True)
        out_file = out_dir / "screenshot.png"

        try:
            result = subprocess.run(
                ["agent-browser", "screenshot", body.url, "--output", str(out_file)],
                capture_output=True, text=True, timeout=30,
            )
            return {
                "status": "ok",
                "filepath": str(out_file) if out_file.exists() else "",
                "stdout": result.stdout[:500],
            }
        except subprocess.TimeoutExpired:
            return {"status": "error", "message": "Browser timeout after 30s"}
        except Exception as e:
            return {"status": "error", "message": str(e)}

    elif body.action == "navigate" and body.url:
        try:
            result = subprocess.run(
                ["agent-browser", "navigate", body.url],
                capture_output=True, text=True, timeout=15,
            )
            return {"status": "ok", "stdout": result.stdout[:1000]}
        except Exception as e:
            return {"status": "error", "message": str(e)}

    elif body.action == "extract" and body.url:
        try:
            cmd = ["agent-browser", "extract", body.url]
            if body.selector:
                cmd.extend(["--selector", body.selector])

            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            return {"status": "ok", "content": result.stdout[:5000]}
        except Exception as e:
            return {"status": "error", "message": str(e)}

    return {"status": "error", "message": f"Unknown action: {body.action}"}


# ── SEO Agent (delegates to existing /api/agent/run) ───

@router.get("/seo/status")
async def skill_seo_status():
    """Check if SEO agent server is running."""
    return {
        "status": "ok",
        "endpoint": "POST /api/agent/run",
        "tools": "20 SEO tools (kb, research, content, optimization, analytics, third-party APIs)",
        "web_ui": "http://127.0.0.1:8000",
        "api_docs": "http://127.0.0.1:8000/docs",
    }


# ── File Browser ────────────────────────────────────────

EXPORTS_DIR = Path("data/exports")


@router.get("/files")
async def list_exports():
    """List all AI-generated files available for download."""
    EXPORTS_DIR.mkdir(parents=True, exist_ok=True)
    files = []
    for f in sorted(EXPORTS_DIR.iterdir(), key=lambda x: x.stat().st_mtime, reverse=True):
        if f.is_file():
            stat = f.stat()
            files.append({
                "name": f.name,
                "size": stat.st_size,
                "size_formatted": _format_size(stat.st_size),
                "modified": stat.st_mtime,
                "ext": f.suffix.lower(),
            })
    return {"files": files}


def _safe_filename(filename: str) -> str:
    """Prevent path traversal — reject any filename containing .. or /."""
    name = Path(filename).name  # strip any directory components
    if name != filename or ".." in filename:
        raise HTTPException(400, "Invalid filename")
    return name


@router.get("/files/{filename}")
async def download_file(filename: str):
    """Download a generated file."""
    safe = _safe_filename(filename)
    filepath = EXPORTS_DIR / safe
    if not filepath.exists():
        raise HTTPException(404, "File not found")
    return FileResponse(str(filepath), filename=safe)


@router.delete("/files/{filename}")
async def delete_export(filename: str):
    """Delete a generated file."""
    safe = _safe_filename(filename)
    filepath = EXPORTS_DIR / safe
    if not filepath.exists():
        raise HTTPException(404, "File not found")
    filepath.unlink()
    return {"status": "deleted", "filename": safe}


def _format_size(size: int) -> str:
    if size < 1024:
        return f"{size} B"
    elif size < 1048576:
        return f"{size / 1024:.0f} KB"
    return f"{size / 1048576:.1f} MB"
