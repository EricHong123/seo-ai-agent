"""Auto-save agent outputs as downloadable files — all four formats."""

import re
import io
from pathlib import Path
from datetime import datetime, timezone

EXPORTS_DIR = Path("data/exports")
EXPORTS_DIR.mkdir(parents=True, exist_ok=True)


def _safe_prefix(prefix: str) -> str:
    return re.sub(r"[^\w\-一-鿿]", "-", prefix)[:40].strip("-")


def _timestamp() -> str:
    return datetime.now(timezone.utc).strftime("%Y%m%d-%H%M%S")


def _parse_sections(text: str) -> list[dict]:
    """Parse markdown into sections: title + body lines."""
    sections = []
    lines = text.strip().split("\n")
    current_title = ""
    current_lines = []

    for line in lines:
        if line.startswith("# ") and not current_title:
            current_title = line[2:].strip()
        elif (line.startswith("## ") or line.startswith("### ")) and current_title:
            if current_lines:
                sections.append({"title": current_title, "body": current_lines})
            current_title = line.lstrip("#").strip()
            current_lines = []
        elif line.startswith("# ") and current_title:
            sections.append({"title": current_title, "body": current_lines})
            current_title = line[2:].strip()
            current_lines = []
        else:
            current_lines.append(line)

    if current_title or current_lines:
        sections.append({"title": current_title or "Report", "body": current_lines})

    return sections


def _parse_table_rows(text: str) -> list[list[str]]:
    """Extract markdown table rows as list of lists."""
    rows = []
    in_table = False
    for line in text.strip().split("\n"):
        line = line.strip()
        if line.startswith("|") and line.endswith("|"):
            cells = [c.strip() for c in line.split("|")[1:-1]]
            # Skip separator rows like |---|----|
            if all(re.match(r"^[-:]+$", c) for c in cells):
                continue
            rows.append(cells)
            in_table = True
        else:
            if in_table and rows:
                break  # end of table
    return rows


# ── MD ────────────────────────────────────────────────

def save_md(content: str, prefix: str) -> str:
    safe = _safe_prefix(prefix)
    ts = _timestamp()
    filename = f"{safe}-{ts}.md"
    (EXPORTS_DIR / filename).write_text(content.strip() + "\n", encoding="utf-8")
    return filename


# ── DOCX ──────────────────────────────────────────────

def save_docx(content: str, prefix: str) -> str:
    try:
        from docx import Document
        from docx.shared import Pt, Inches
    except ImportError:
        return ""

    doc = Document()
    sections = _parse_sections(content)

    for i, sec in enumerate(sections):
        if i > 0:
            doc.add_page_break()
        doc.add_heading(sec["title"], level=1)
        for line in sec["body"]:
            line = line.strip()
            if not line:
                continue
            if line.startswith("|") and "|" in line[1:]:
                # Table row — handled below
                continue
            if line.startswith("- ") or line.startswith("* "):
                doc.add_paragraph(line[2:], style="List Bullet")
            elif re.match(r"^\d+[.)]\s", line):
                doc.add_paragraph(re.sub(r"^\d+[.)]\s*", "", line), style="List Number")
            elif line.startswith("**") and "**" in line[2:]:
                p = doc.add_paragraph()
                run = p.add_run(line.strip("*"))
                run.bold = True
            else:
                doc.add_paragraph(line)

    safe = _safe_prefix(prefix)
    ts = _timestamp()
    filename = f"{safe}-{ts}.docx"
    doc.save(str(EXPORTS_DIR / filename))
    return filename


# ── PPTX ──────────────────────────────────────────────

def save_pptx(content: str, prefix: str) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return ""

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    sections = _parse_sections(content)
    if not sections:
        sections = [{"title": prefix or "Report", "body": content.strip().split("\n")}]

    for sec in sections:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = sec["title"][:100]

        body_lines = [l for l in sec["body"] if l.strip()]
        if body_lines:
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            for bl in body_lines[:20]:
                p = tf.add_paragraph()
                p.text = bl[:200]

    safe = _safe_prefix(prefix)
    ts = _timestamp()
    filename = f"{safe}-{ts}.pptx"
    prs.save(str(EXPORTS_DIR / filename))
    return filename


# ── XLSX ──────────────────────────────────────────────

def save_xlsx(content: str, prefix: str) -> str:
    try:
        import openpyxl
    except ImportError:
        return ""

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"

    # Try to extract table first
    rows = _parse_table_rows(content)

    if rows:
        for row_idx, cells in enumerate(rows, 1):
            for col_idx, val in enumerate(cells, 1):
                ws.cell(row=row_idx, column=col_idx, value=val)
    else:
        # No table found — put key-value pairs or full text
        lines = [l.strip() for l in content.strip().split("\n") if l.strip()]
        ws.cell(row=1, column=1, value="Content")
        for i, line in enumerate(lines[:200], 2):
            ws.cell(row=i, column=1, value=line[:500])

    safe = _safe_prefix(prefix)
    ts = _timestamp()
    filename = f"{safe}-{ts}.xlsx"
    wb.save(str(EXPORTS_DIR / filename))
    return filename


# ── Unified export ─────────────────────────────────────

def save_all_formats(content: str, prefix: str = "output") -> dict[str, str]:
    """Save content as MD, DOCX, PPTX, XLSX. Returns {ext: filename}."""
    result = {}
    result["md"] = save_md(content, prefix)
    result["docx"] = save_docx(content, prefix)
    result["pptx"] = save_pptx(content, prefix)
    result["xlsx"] = save_xlsx(content, prefix)
    return result
