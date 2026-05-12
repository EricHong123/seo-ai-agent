"""PPTX generation tool — agent can call this to create PowerPoint files."""

from llm.base import ToolDef


TOOL_SCHEMA = {
    "name": "generate_pptx",
    "description": "Generate a PowerPoint (.pptx) presentation file. Input Markdown content — use '---' to separate slides, '# Title' for slide titles. The file will be saved and available for download. Use this when the user asks to create a presentation, slide deck, or PPT.",
    "parameters": {
        "type": "object",
        "properties": {
            "title": {"type": "string", "description": "Presentation title / filename"},
            "content": {"type": "string", "description": "Markdown content for slides. Separate slides with '---'. Use '# Slide Title' for each slide heading."},
        },
        "required": ["title", "content"],
    },
}


def make_tool() -> ToolDef:
    async def handler(title: str, content: str) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from pathlib import Path
        except ImportError:
            return "Error: python-pptx not installed. Run: pip install python-pptx"

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slides_text = content.split("\n---\n")
        slide_count = 0

        for slide_md in slides_text:
            lines = slide_md.strip().split("\n")
            title_text = ""
            body_lines = []

            for line in lines:
                if (line.startswith("# ") or line.startswith("## ")) and not title_text:
                    title_text = line.lstrip("#").strip()
                else:
                    body_lines.append(line)

            if not title_text and body_lines:
                title_text = body_lines[0][:100]
                body_lines = body_lines[1:]

            if not title_text:
                title_text = f"Slide {slide_count + 1}"

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_text[:100]

            if body_lines:
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                for bl in body_lines[:20]:
                    p = tf.add_paragraph()
                    p.text = bl[:200]

            slide_count += 1

        out_dir = Path("data/exports")
        out_dir.mkdir(parents=True, exist_ok=True)
        filename = title.replace(" ", "-").replace("/", "-")[:40]
        out_path = out_dir / f"{filename}.pptx"
        prs.save(str(out_path))

        return (
            f"PPT generated: {out_path.name} ({slide_count} slides). "
            f"Download: /api/skills/files/{out_path.name}"
        )

    return ToolDef(
        name=TOOL_SCHEMA["name"],
        description=TOOL_SCHEMA["description"],
        parameters=TOOL_SCHEMA["parameters"],
        handler=handler,
    )
